from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation

# Windows 默认 stdout 是 cp936/gbk，遇到 −/✓/Δ 等 Unicode 字符会 UnicodeEncodeError。
# 这里强制切到 UTF-8，errors=replace 兜底以避免脚本因终端编码崩溃。
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def shape_texts(shape) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text:
        text = " ".join(shape.text.split())
        if text:
            texts.append(text)
    if getattr(shape, "has_table", False):
        for cell in shape.table.iter_cells():
            text = " ".join(cell.text.split())
            if text:
                texts.append(text)
    return texts


def main() -> int:
    parser = argparse.ArgumentParser(description="Scan PPTX text and optional forbidden terms.")
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--bad", default="", help="Comma-separated terms that should not remain")
    parser.add_argument("--json", default="", help="Optional JSON output path")
    args = parser.parse_args()

    pptx = Path(args.pptx).expanduser().resolve()
    bad_terms = [x.strip() for x in args.bad.split(",") if x.strip()]
    prs = Presentation(str(pptx))

    slides = []
    hits = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in iter_shapes(slide.shapes):
            texts.extend(shape_texts(shape))
        joined = "\n".join(texts)
        slide_hits = [term for term in bad_terms if term in joined]
        if slide_hits:
            hits.append({"slide": idx, "terms": slide_hits})
        slides.append({"slide": idx, "text": joined})

    payload = {"pptx": str(pptx), "slide_count": len(prs.slides), "bad_hits": hits, "slides": slides}
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    if args.json:
        out = Path(args.json).expanduser().resolve()
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(text, encoding="utf-8")
    print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
