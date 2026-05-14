"""Dump every slide's shape/text/table/picture inventory to a Markdown file.

迭代填充时强烈推荐先跑这个脚本拿到模板/骨架的真实文本,再去构造
``replace_exact_text`` 的 dict 或 ``write_table`` 的二维数组。直接靠肉眼
猜模板里的字符串很容易因为标点/空格/引号种类不一致而漏替换。

Usage::

    python dump_pptx_content.py --pptx deck.pptx --output dump.md

输出会按 slide 编号分节,每节列出:

- 每个 shape 的索引/名称/类型/文本 (text_frame)
- 每个 table 的二维内容
- 每张 picture 的 shape.name + 尺寸 + 位置

文件用 UTF-8 写盘,可以直接交给后续脚本做精确替换。
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# Windows 默认 stdout 是 cp936/gbk,中文/特殊字符易触发 UnicodeEncodeError。
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def fmt_emu(value) -> str:
    try:
        return f"{Emu(value).inches:.2f}in"
    except Exception:
        return str(value)


def dump_shape(shape, depth: int = 0) -> list[str]:
    indent = "  " * depth
    lines: list[str] = []
    name = getattr(shape, "name", "?")
    stype = getattr(shape, "shape_type", "?")
    pos = ""
    try:
        pos = (
            f" @ left={fmt_emu(shape.left)} top={fmt_emu(shape.top)}"
            f" w={fmt_emu(shape.width)} h={fmt_emu(shape.height)}"
        )
    except Exception:
        pass

    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        preview = text.replace("\n", " / ")
        if len(preview) > 240:
            preview = preview[:237] + "..."
        lines.append(f"{indent}- TEXT [{stype}] `{name}`{pos}")
        if text:
            lines.append(f"{indent}  > {preview}")
        else:
            lines.append(f"{indent}  > (empty)")
    elif getattr(shape, "has_table", False):
        lines.append(f"{indent}- TABLE `{name}`{pos}")
        for r_idx, row in enumerate(shape.table.rows):
            cells = [
                c.text_frame.text.strip().replace("\n", " ⏎ ")
                for c in row.cells
            ]
            lines.append(f"{indent}  | r{r_idx} | " + " | ".join(cells) + " |")
    elif int(stype) == 13:  # PICTURE
        lines.append(f"{indent}- PICTURE `{name}`{pos}")
    elif int(stype) == 6:  # GROUP
        lines.append(f"{indent}- GROUP `{name}`{pos}")
        for child in shape.shapes:
            lines.extend(dump_shape(child, depth + 1))
    else:
        lines.append(f"{indent}- SHAPE [{stype}] `{name}`{pos}")
    return lines


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Dump every slide's shapes/text/tables/pictures to a Markdown file."
    )
    parser.add_argument("--pptx", required=True, help="Input PPTX path")
    parser.add_argument("--output", required=True, help="Output Markdown path")
    parser.add_argument(
        "--slide",
        default="",
        help="Comma-separated 1-based slide indices to dump (default: all)",
    )
    args = parser.parse_args()

    pptx_path = Path(args.pptx).expanduser().resolve()
    out_path = Path(args.output).expanduser().resolve()
    if not pptx_path.exists():
        raise SystemExit(f"PPTX not found: {pptx_path}")

    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    if args.slide.strip():
        wanted = {int(x.strip()) for x in args.slide.split(",") if x.strip()}
    else:
        wanted = set(range(1, total + 1))

    out_lines: list[str] = []
    out_lines.append(f"# PPTX Content Dump: {pptx_path.name}")
    out_lines.append("")
    out_lines.append(f"- Source: `{pptx_path}`")
    out_lines.append(f"- Total slides: {total}")
    sw = fmt_emu(prs.slide_width)
    sh = fmt_emu(prs.slide_height)
    out_lines.append(f"- Slide size: {sw} x {sh}")
    out_lines.append("")

    for idx, slide in enumerate(prs.slides, 1):
        if idx not in wanted:
            continue
        layout_name = getattr(slide.slide_layout, "name", "?")
        out_lines.append(f"## Slide {idx} (layout: {layout_name})")
        out_lines.append("")
        # 注意:这里只遍历直接子 shapes,group 内部由 dump_shape 递归处理,
        # 保留缩进层级,便于看清模板嵌套结构。
        for shape in slide.shapes:
            out_lines.extend(dump_shape(shape, depth=0))
        out_lines.append("")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(out_lines), encoding="utf-8")
    print(out_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
