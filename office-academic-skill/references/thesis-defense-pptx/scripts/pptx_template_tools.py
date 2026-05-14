from __future__ import annotations

import os
from pathlib import Path
from typing import Sequence

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def _brand_red() -> RGBColor:
    raw = os.getenv("THESIS_PPTX_BRAND_RED", "").strip()
    if raw:
        try:
            if raw.startswith("#"):
                raw = raw[1:]
            if len(raw) == 6 and all(c in "0123456789abcdefABCDEF" for c in raw):
                return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
            parts = [int(x.strip()) for x in raw.split(",")]
            if len(parts) == 3 and all(0 <= x <= 255 for x in parts):
                return RGBColor(parts[0], parts[1], parts[2])
        except ValueError:
            pass
    return RGBColor(0x83, 0x06, 0x04)


RED = _brand_red()
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)


def emu(inches: float) -> int:
    return Inches(inches)


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def set_text_preserve_style(shape, text: str, size: float | None = None, bold: bool | None = None, align=None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    old_size = None
    old_bold = None
    old_color = None
    old_name = None
    for old_p in tf.paragraphs:
        for old_r in old_p.runs:
            if old_r.text.strip():
                old_size = old_r.font.size
                old_bold = old_r.font.bold
                old_name = old_r.font.name
                try:
                    old_color = old_r.font.color.rgb
                except Exception:
                    old_color = None
                break
        if old_size is not None or old_color is not None:
            break

    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    if old_name:
        r.font.name = old_name
    if size is not None:
        r.font.size = Pt(size)
    elif old_size is not None:
        r.font.size = old_size
    if bold is not None:
        r.font.bold = bold
    elif old_bold is not None:
        r.font.bold = old_bold
    if old_color is not None:
        r.font.color.rgb = old_color


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: float, bold: bool = False, color=BLACK, align=None):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = emu(0.02)
    tf.margin_bottom = emu(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align or PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_para(slide, lines: Sequence[str], x: float, y: float, w: float, h: float, size: float, color=BLACK, bold_first: bool = False):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = emu(0.02)
    tf.margin_bottom = emu(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bool(bold_first and i == 0)
        r.font.color.rgb = color
    return box


def rect(slide, x: float, y: float, w: float, h: float, line=RED, fill=None, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, emu(x), emu(y), emu(w), emu(h))
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape


def tag(slide, text: str, x: float, y: float, w: float, h: float, fill=RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    set_text_preserve_style(shape, text, 20, False, PP_ALIGN.CENTER)
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = WHITE
    return shape


def add_pic_fit(slide, src: str | Path, x: float, y: float, w: float, h: float):
    path = Path(src)
    with Image.open(path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    box_ratio = w / h
    if img_ratio > box_ratio:
        nw = w
        nh = w / img_ratio
    else:
        nh = h
        nw = h * img_ratio
    return slide.shapes.add_picture(str(path), emu(x + (w - nw) / 2), emu(y + (h - nh) / 2), width=emu(nw), height=emu(nh))


def replace_exact_text(slide, mapping: dict[str, str]) -> None:
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text in mapping:
                set_text_preserve_style(shape, mapping[text])


def replace_partial_text(slide, mapping: dict[str, str], min_len: int = 0) -> int:
    """按 substring 匹配批量替换文本。

    与 ``replace_exact_text`` 互补:模板里小标点差异(如 ``——`` vs ``—``、空格、
    全/半角引号)很容易让等值匹配失败，此处用包含匹配兜底。

    Args:
        mapping: ``{substring: new_full_text}``。命中任一 key 的 shape，
            其整段文字将被替换为对应的 ``new_full_text``。
        min_len: 仅替换原文长度不小于此值的 shape，避免短 key (如"状态")
            误伤导航标签。默认 0 不过滤。

    匹配顺序按 key 长度从长到短，避免短 key 提前命中。返回替换次数。
    """
    keys = sorted(mapping.keys(), key=len, reverse=True)
    n = 0
    for shape in iter_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if len(text) < min_len:
            continue
        for key in keys:
            if key and key in text:
                set_text_preserve_style(shape, mapping[key])
                n += 1
                break
    return n


def find_shape_by_text(slide, needle: str):
    """返回第一个文字包含 ``needle`` 的 shape，找不到返回 ``None``。"""
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False) and needle in shape.text:
            return shape
    return None


def find_picture(slide, name_substr: str | None = None):
    """返回 slide 中第一张图片 shape；可选 ``name_substr`` 过滤 shape.name。"""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            if name_substr is None or name_substr in shape.name:
                return shape
    return None


def find_table(slide):
    """返回 slide 中第一个 table 对象，找不到返回 ``None``。"""
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return shape.table
    return None


def replace_picture(slide, old_pic, new_path: str | Path):
    """用 ``new_path`` 指向的新图替换原 picture shape，保持 left/top/width/height。

    python-pptx 没有原地换图 API，标准做法是删除旧 shape 后在原位置 add_picture。
    返回新插入的 picture shape。
    """
    left, top, width, height = old_pic.left, old_pic.top, old_pic.width, old_pic.height
    sp = old_pic._element
    sp.getparent().remove(sp)
    return slide.shapes.add_picture(str(Path(new_path)), left, top, width=width, height=height)


def write_table(table, rows: Sequence[Sequence[str]], font_size: float | None = None) -> None:
    """按二维数组覆盖表格单元格文字，保留每个单元格首个 run 的样式。

    适合在模板里"先有表格框架再换内容"的场景。新写入的内容若多于原表格
    维度，超出部分会被忽略，不会自动加行/加列。

    Args:
        table: ``shape.table`` 对象。
        rows: 二维序列，``rows[r][c]`` 写入到 ``table.cell(r, c)``。
        font_size: 显式指定字号(磅)；为 ``None`` 时沿用原单元格字号。
    """
    for r_idx, row in enumerate(rows):
        if r_idx >= len(table.rows):
            break
        for c_idx, cell_text in enumerate(row):
            if c_idx >= len(table.columns):
                break
            cell = table.cell(r_idx, c_idx)
            old_size = None
            old_bold = None
            old_color = None
            old_name = None
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    if run.text.strip():
                        old_size = run.font.size
                        old_bold = run.font.bold
                        old_name = run.font.name
                        try:
                            old_color = run.font.color.rgb
                        except Exception:
                            old_color = None
                        break
                if old_size is not None:
                    break
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            if old_name:
                run.font.name = old_name
            if old_bold is not None:
                run.font.bold = old_bold
            if old_color is not None:
                run.font.color.rgb = old_color
            if font_size is not None:
                run.font.size = Pt(font_size)
            elif old_size is not None:
                run.font.size = old_size
