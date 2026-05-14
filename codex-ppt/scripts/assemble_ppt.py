#!/usr/bin/env python3
"""
PPT 组装脚本

将一个目录中的幻灯片图片组装成 PowerPoint 演示文稿。
每张图片会被插入为一页幻灯片，充满整个页面。
"""

import argparse
import os
import re
import sys
import tempfile
from typing import Dict, List, Optional


def dependency_hint() -> str:
    runtime_home = os.path.expanduser(os.environ.get("CODEX_PPT_HOME", "~/.codex-ppt-skill"))
    python = os.path.join(
        runtime_home,
        ".venv",
        "Scripts" if os.name == "nt" else "bin",
        "python.exe" if os.name == "nt" else "python",
    )
    runtime_script = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "codex_ppt_runtime.py",
    )
    return (
        f"请运行: python3 {runtime_script} bootstrap\n"
        f"或直接运行: {python} -m pip install -r "
        f"{os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'requirements.txt')}"
    )


def get_slide_images(ppt_project_dir: str) -> List[str]:
    """
    获取幻灯片图片文件列表，按文件名排序

    Args:
        ppt_project_dir: PPT 项目目录（包含 origin_image 子目录）

    Returns:
        按顺序排列的图片文件路径列表
    """
    # 从 origin_image 子目录读取
    origin_image_dir = os.path.join(ppt_project_dir, "origin_image")

    if not os.path.exists(origin_image_dir):
        print(f"错误：origin_image 目录不存在: {origin_image_dir}")
        return []

    print(f"从 origin_image 目录读取图片: {origin_image_dir}")

    # 支持的图片格式
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp'}

    # 只获取正式幻灯片图片，避免 sample_slide.png、草稿图或参考图被误装入 PPT。
    slide_name_pattern = re.compile(r"^slide_(\d+)\.(png|jpe?g|gif|bmp)$", re.IGNORECASE)
    image_files = []
    for file in os.listdir(origin_image_dir):
        file_path = os.path.join(origin_image_dir, file)
        if os.path.isfile(file_path):
            ext = os.path.splitext(file)[1].lower()
            if ext in image_extensions and slide_name_pattern.match(file):
                image_files.append(file_path)

    # 按文件名排序
    image_files.sort()

    return image_files


def load_speaker_notes(ppt_project_dir: str) -> Dict[int, str]:
    """
    从 speech.md 读取每页演讲备注。

    支持以下标题格式：
    - ## Slide 1: 标题
    - ### Slide 1: 标题
    - ## 第 1 页：标题
    """
    speech_path = os.path.join(ppt_project_dir, "speech.md")
    if not os.path.exists(speech_path):
        return {}

    with open(speech_path, "r", encoding="utf-8") as f:
        content = f.read()

    notes: Dict[int, str] = {}
    current_slide: Optional[int] = None
    current_lines: List[str] = []
    heading_pattern = re.compile(r"^#{2,4}\s*(?:Slide\s*(\d+)|第\s*(\d+)\s*页)\b.*$", re.IGNORECASE)

    def flush_current() -> None:
        if current_slide is None:
            return
        text = "\n".join(current_lines).strip()
        if text:
            notes[current_slide] = text

    for line in content.splitlines():
        match = heading_pattern.match(line.strip())
        if match:
            flush_current()
            current_slide = int(match.group(1) or match.group(2))
            current_lines = []
            continue

        if current_slide is not None:
            current_lines.append(line)

    flush_current()
    return notes


def compress_image_if_needed(
    image_path: str,
    max_size_mb: float = 2.0,
    quality_step: int = 5
) -> Optional[str]:
    """
    如果图片超过指定大小，压缩图片并返回临时文件路径

    Args:
        image_path: 原始图片路径
        max_size_mb: 最大文件大小（MB）
        quality_step: 每次降低的质量步进

    Returns:
        str: 如果需要压缩，返回临时文件路径；否则返回 None
    """
    max_size_bytes = max_size_mb * 1024 * 1024

    # 检查原始文件大小
    file_size = os.path.getsize(image_path)

    if file_size <= max_size_bytes:
        # 不需要压缩
        return None

    print(f"  图片大小 {file_size / 1024 / 1024:.2f}MB，需要压缩...")

    try:
        from PIL import Image

        # 打开图片
        img = Image.open(image_path)

        # 转换 RGBA 到 RGB（如果需要保存为 JPEG）
        if img.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
            img = background

        # 创建临时文件
        temp_fd, temp_path = tempfile.mkstemp(suffix='.jpg')
        os.close(temp_fd)

        # 从高质量开始尝试压缩
        quality = 95
        while quality > 20:
            img.save(temp_path, 'JPEG', quality=quality, optimize=True)
            compressed_size = os.path.getsize(temp_path)

            if compressed_size <= max_size_bytes:
                print(f"  压缩成功: {compressed_size / 1024 / 1024:.2f}MB (质量: {quality})")
                return temp_path

            quality -= quality_step

        # 如果还是太大，尝试调整尺寸
        print(f"  质量压缩不足，尝试缩小尺寸...")
        scale = 0.9
        while scale > 0.3:
            new_width = int(img.width * scale)
            new_height = int(img.height * scale)
            resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

            resized_img.save(temp_path, 'JPEG', quality=85, optimize=True)
            compressed_size = os.path.getsize(temp_path)

            if compressed_size <= max_size_bytes:
                print(f"  压缩成功: {compressed_size / 1024 / 1024:.2f}MB (缩放: {scale:.0%})")
                return temp_path

            scale -= 0.1

        # 实在压不下去了，返回最后的结果
        print(f"  警告：无法压缩到 {max_size_mb}MB 以下，使用最小尺寸版本")
        return temp_path

    except ImportError:
        print("错误：未安装 Pillow 库")
        print(dependency_hint())
        return None
    except Exception as e:
        print(f"  警告：图片压缩失败: {e}")
        if 'temp_path' in locals() and os.path.exists(temp_path):
            os.remove(temp_path)
        return None


def create_presentation(
    image_files: List[str],
    output_path: str,
    aspect_ratio: str = "16:9",
    speaker_notes: Optional[Dict[int, str]] = None,
) -> bool:
    """
    创建 PowerPoint 演示文稿

    Args:
        image_files: 幻灯片图片文件列表
        output_path: 输出 PPT 文件路径
        aspect_ratio: 幻灯片宽高比（16:9 或 4:3）

    Returns:
        bool: 成功返回 True，失败返回 False
    """
    try:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            print("错误：未安装 python-pptx 库")
            print(dependency_hint())
            return False

        # 创建演示文稿
        prs = Presentation()

        # 设置幻灯片尺寸
        if aspect_ratio == "16:9":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
        elif aspect_ratio == "4:3":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        else:
            print(f"警告：不支持的宽高比 {aspect_ratio}，使用默认值 16:9")
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

        speaker_notes = speaker_notes or {}

        # 添加每一页幻灯片
        temp_files_to_cleanup = []
        for i, image_path in enumerate(image_files, 1):
            if not os.path.exists(image_path):
                print(f"警告：图片文件不存在: {image_path}")
                continue

            # 压缩图片（如果需要）
            compressed_path = compress_image_if_needed(image_path, max_size_mb=2.0)

            # 使用压缩后的图片或原图
            image_to_use = compressed_path if compressed_path else image_path

            # 记录临时文件以便后续清理
            if compressed_path:
                temp_files_to_cleanup.append(compressed_path)

            # 使用空白布局（索引 6）
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 将图片添加到幻灯片，充满整个页面
            slide.shapes.add_picture(
                image_to_use,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )

            note_text = speaker_notes.get(i)
            if note_text:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.clear()
                notes_frame.text = note_text

            print(f"✓ 已添加第 {i} 页: {os.path.basename(image_path)}")

        # 确保输出目录存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # 保存演示文稿
        prs.save(output_path)
        print(f"\n✓ PPT 文件已保存: {output_path}")
        print(f"  总页数: {len(image_files)}")
        if speaker_notes:
            matched_notes = sum(1 for i in range(1, len(image_files) + 1) if speaker_notes.get(i))
            print(f"  已写入备注: {matched_notes}/{len(image_files)} 页")

        # 清理临时文件
        for temp_file in temp_files_to_cleanup:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                print(f"警告：清理临时文件失败 {temp_file}: {e}")

        return True

    except Exception as e:
        print(f"错误：创建 PPT 失败: {e}")
        import traceback
        traceback.print_exc()

        # 即使失败也要清理临时文件
        if 'temp_files_to_cleanup' in locals():
            for temp_file in temp_files_to_cleanup:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except:
                    pass

        return False


def main():
    parser = argparse.ArgumentParser(
        description='将幻灯片图片组装成 PowerPoint 演示文稿',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
示例用法:
  # 基本用法
  # 会在 /path/to/base/ 下查找 MyPresentation/ 文件夹
  # 从 MyPresentation/origin_image/ 读取图片
  # 将 PPT 保存为 MyPresentation/MyPresentation.pptx
  python assemble_ppt.py /path/to/base/ MyPresentation.pptx

  # 指定 4:3 宽高比
  python assemble_ppt.py /path/to/base/ MyPresentation.pptx --aspect-ratio 4:3

  # 只初始化目录，不生成 PPT
  python assemble_ppt.py /path/to/base/ MyPresentation.pptx --init

文件夹结构要求:
  /path/to/base/MyPresentation/
  ├── origin_image/
  │   ├── slide_01.png
  │   ├── slide_02.png
  │   └── ...
  ├── speech.md (可选，将写入 PPT 备注)
  └── MyPresentation.pptx (将在此生成)

注意：
  - 图片文件必须放在 origin_image 子目录中
  - 只会读取 slide_01.png、slide_02.png 这类正式图片，其他图片会被忽略
  - 图片文件按文件名排序
  - 建议图片文件命名为: slide_01.png, slide_02.png, ...
  - 每张图片会充满整个幻灯片页面
  - 如果项目目录下存在 speech.md，会按 Slide N 标题写入每页备注
        '''
    )

    parser.add_argument('base_dir', help='基础目录（PPT 项目文件夹的父目录）')
    parser.add_argument('output', help='输出 PPT 文件名 (.pptx)')
    parser.add_argument('--aspect-ratio', '--ar',
                        choices=['16:9', '4:3'],
                        default='16:9',
                        help='幻灯片宽高比 (默认: 16:9)')
    parser.add_argument('--init',
                        action='store_true',
                        help='只创建 PPT 项目目录和 origin_image 子目录，不生成 PPT')

    args = parser.parse_args()

    # 确保输出文件有 .pptx 扩展名
    output_filename = args.output
    if not output_filename.lower().endswith('.pptx'):
        output_filename += '.pptx'

    # 获取 PPT 名称（不含扩展名）
    ppt_name = os.path.splitext(os.path.basename(output_filename))[0]

    # 构建 PPT 项目目录
    ppt_project_dir = os.path.join(args.base_dir, ppt_name)
    origin_image_dir = os.path.join(ppt_project_dir, "origin_image")

    if args.init:
        os.makedirs(origin_image_dir, exist_ok=True)
        print(f"✓ PPT 项目目录已准备好: {ppt_project_dir}")
        print(f"✓ 幻灯片图片目录已准备好: {origin_image_dir}")
        sys.exit(0)

    if not os.path.exists(ppt_project_dir):
        print(f"错误：PPT 项目目录不存在: {ppt_project_dir}")
        print("如需初始化目录，请添加 --init 参数")
        sys.exit(1)

    if not os.path.exists(origin_image_dir):
        print(f"错误：origin_image 目录不存在: {origin_image_dir}")
        print("如需初始化目录，请添加 --init 参数")
        sys.exit(1)

    # 设置输出路径
    output_path = os.path.join(ppt_project_dir, output_filename)

    # 获取幻灯片图片
    print(f"正在扫描 PPT 项目目录: {ppt_project_dir}")
    image_files = get_slide_images(ppt_project_dir)

    if not image_files:
        print("错误：未找到任何图片文件")
        print("支持的格式: .png, .jpg, .jpeg, .gif, .bmp")
        print(f"\n请将幻灯片图片放置在: {origin_image_dir}/")
        sys.exit(1)

    print(f"找到 {len(image_files)} 张幻灯片图片\n")
    speaker_notes = load_speaker_notes(ppt_project_dir)
    if speaker_notes:
        print(f"找到 {len(speaker_notes)} 页备注: {os.path.join(ppt_project_dir, 'speech.md')}\n")

    # 创建演示文稿
    print(f"正在创建 PPT (宽高比: {args.aspect_ratio})...")
    print("-" * 50)

    success = create_presentation(image_files, output_path, args.aspect_ratio, speaker_notes)

    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
